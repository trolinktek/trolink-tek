from pathlib import Path
import json
from pypdf import PdfReader
from pptx import Presentation
from openpyxl import load_workbook

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "source-materials"
OUT = ROOT / "research"
OUT.mkdir(exist_ok=True)

report = {"documents": {}, "workbook": {}}

for path in sorted(SRC.iterdir()):
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        reader = PdfReader(path)
        pages = []
        for idx, page in enumerate(reader.pages, 1):
            text = (page.extract_text() or "").strip()
            pages.append({"page": idx, "text": text})
        report["documents"][path.name] = {"type": "pdf", "pages": pages}
    elif suffix == ".pptx":
        prs = Presentation(path)
        slides = []
        for idx, slide in enumerate(prs.slides, 1):
            chunks = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    chunks.append(shape.text.strip())
            slides.append({"slide": idx, "text": "\n".join(chunks)})
        report["documents"][path.name] = {"type": "pptx", "slides": slides}
    elif suffix == ".xlsx":
        wb = load_workbook(path, data_only=False, read_only=True)
        sheets = {}
        for ws in wb.worksheets:
            rows = []
            for row in ws.iter_rows(values_only=True):
                vals = [v for v in row]
                if any(v not in (None, "") for v in vals):
                    rows.append(vals)
            sheets[ws.title] = rows
        report["workbook"] = {"file": path.name, "sheets": sheets}

(OUT / "source-extraction.json").write_text(
    json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8"
)

summary = []
for name, doc in report["documents"].items():
    summary.append(f"# {name}")
    units = doc.get("pages", doc.get("slides", []))
    key = "page" if doc["type"] == "pdf" else "slide"
    for unit in units:
        if unit["text"]:
            summary.append(f"\n## {key.title()} {unit[key]}\n\n{unit['text']}")
for sheet, rows in report.get("workbook", {}).get("sheets", {}).items():
    summary.append(f"\n# Workbook sheet: {sheet}")
    for row in rows:
        summary.append(" | ".join("" if v is None else str(v) for v in row))
(OUT / "source-extraction.md").write_text("\n".join(summary), encoding="utf-8")

print(json.dumps({
    "documents": {k: len(v.get("pages", v.get("slides", []))) for k, v in report["documents"].items()},
    "sheets": {k: len(v) for k, v in report.get("workbook", {}).get("sheets", {}).items()},
}, ensure_ascii=False, indent=2))
